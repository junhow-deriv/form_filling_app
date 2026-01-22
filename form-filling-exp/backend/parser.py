"""
File parsing module using Docling.

Provides functionality to parse various file types (PDF, PPTX, DOCX, images)
into markdown format for use as context in the form-filling agent.
"""

import asyncio
import os
from pathlib import Path
from typing import AsyncGenerator, Literal, Optional

# Maximum concurrent parsing requests
MAX_CONCURRENT_PARSE = 5

# File extensions that don't need parsing (already text-based)
SIMPLE_TEXT_EXTENSIONS = {
    '.txt', '.md', '.markdown', '.csv', '.json', '.xml', '.html', '.htm',
    '.py', '.js', '.ts', '.jsx', '.tsx', '.css', '.scss', '.yaml', '.yml',
    '.toml', '.ini', '.cfg', '.conf', '.sh', '.bash', '.zsh', '.sql',
    '.r', '.rb', '.go', '.java', '.c', '.cpp', '.h', '.hpp', '.rs',
    '.adoc', '.asciidoc', '.vtt', '.xhtml',
}

# File extensions that need Docling
PARSEABLE_EXTENSIONS = {
    '.pdf', '.pptx', '.ppt', '.docx', '.doc', '.xlsx', '.xls',
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp',
}

# Parse mode options (kept for compatibility, though Docling handles this differently)
ParseMode = Literal["cost_effective", "agentic_plus"]

# Try to import Docling
DOCLING_AVAILABLE = False
DOCLING_ERROR = None

try:
    from docling.document_converter import DocumentConverter
    DOCLING_AVAILABLE = True
except ImportError as e:
    DOCLING_ERROR = str(e)
    print(f"[Parser] Docling not available: {e}")
    print("[Parser] Install with: pip install docling")


def needs_parsing(filename: str) -> bool:
    """Check if a file needs to be parsed with Docling."""
    ext = Path(filename).suffix.lower()
    return ext in PARSEABLE_EXTENSIONS


def is_simple_text(filename: str) -> bool:
    """Check if a file is simple text that can be read directly."""
    ext = Path(filename).suffix.lower()
    return ext in SIMPLE_TEXT_EXTENSIONS


def estimate_file_chars(file_bytes: bytes, filename: str) -> int:
    """
    Quickly estimate the character count of a file's text content.

    This is used for pre-validation before expensive parsing calls.
    The estimate is intentionally conservative (may overestimate) to avoid
    wasting time on files that will be rejected anyway.

    Args:
        file_bytes: The file content as bytes
        filename: The original filename

    Returns:
        Estimated character count of extracted text
    """
    ext = Path(filename).suffix.lower()

    # Simple text files - just return byte length (roughly char count for UTF-8)
    if is_simple_text(filename):
        return len(file_bytes)

    # Images - OCR output varies wildly, use conservative fixed estimate
    # A typical page of OCR'd text is ~2000-3000 chars
    if ext in {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'}:
        # Estimate based on file size - larger images often have more text
        # But cap it since very large images might just be high-res photos
        size_mb = len(file_bytes) / (1024 * 1024)
        return min(int(size_mb * 3000), 10000)  # ~3k chars per MB, max 10k

    # PDFs - use PyMuPDF for quick text extraction
    if ext == '.pdf':
        try:
            import fitz
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            total_chars = 0
            for page in doc:
                text = page.get_text()
                total_chars += len(text)
            doc.close()
            # Add 20% buffer since parsing often extracts more (tables, etc.)
            return int(total_chars * 1.2)
        except Exception as e:
            print(f"[Estimate] PDF extraction failed for {filename}: {e}")
            # Fallback: estimate based on file size (~1000 chars per KB for PDFs)
            return len(file_bytes) // 10

    # DOCX - use python-docx for quick text extraction
    if ext in {'.docx', '.doc'}:
        try:
            from docx import Document
            import io
            doc = Document(io.BytesIO(file_bytes))
            total_chars = 0
            for para in doc.paragraphs:
                total_chars += len(para.text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        total_chars += len(cell.text)
            return int(total_chars * 1.1)  # Small buffer
        except ImportError:
            print("[Estimate] python-docx not available, using file size heuristic")
            return len(file_bytes) // 5  # ~200 chars per KB for DOCX
        except Exception as e:
            print(f"[Estimate] DOCX extraction failed for {filename}: {e}")
            return len(file_bytes) // 5

    # PPTX - use python-pptx for quick text extraction
    if ext in {'.pptx', '.ppt'}:
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(file_bytes))
            total_chars = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        total_chars += len(shape.text)
            return int(total_chars * 1.2)  # Buffer for tables/images
        except Exception as e:
            print(f"[Estimate] PPTX extraction failed for {filename}: {e}")
            return len(file_bytes) // 8  # ~125 chars per KB for PPTX

    # XLSX - estimate based on file size
    if ext in {'.xlsx', '.xls'}:
        # Spreadsheets can be very dense with data
        return len(file_bytes) // 3  # ~333 chars per KB

    # Unknown file type - conservative estimate
    return len(file_bytes) // 4


async def parse_file(
    file_bytes: bytes,
    filename: str,
    mode: ParseMode = "cost_effective",
    api_key: Optional[str] = None,
) -> str:
    """
    Parse a file and return its markdown content using Docling.

    Args:
        file_bytes: The file content as bytes
        filename: The original filename
        mode: The parsing mode (ignored for Docling, kept for compatibility)
        api_key: API key (ignored for Docling, kept for compatibility)

    Returns:
        Markdown string of the parsed content
    """
    import tempfile

    ext = Path(filename).suffix.lower()

    # If it's a simple text file, just decode and return
    if is_simple_text(filename):
        try:
            return file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            return file_bytes.decode('latin-1')

    # If it doesn't need parsing and isn't simple text, return error
    if not needs_parsing(filename):
        return f"[Unsupported file type: {ext}]"

    # Use Docling
    if not DOCLING_AVAILABLE:
        raise RuntimeError(f"Docling not available: {DOCLING_ERROR}")

    # Write bytes to temp file (Docling takes file paths)
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        # Run Docling in a thread since it's synchronous/CPU-bound
        def _run_docling():
            converter = DocumentConverter()
            result = converter.convert(tmp_path)
            return result.document.export_to_markdown()

        return await asyncio.to_thread(_run_docling)
    finally:
        # Clean up temp file
        try:
            os.unlink(tmp_path)
        except Exception:
            pass


async def parse_files_stream(
    files: list[tuple[bytes, str]],
    mode: ParseMode = "cost_effective",
    api_key: Optional[str] = None,
) -> AsyncGenerator[dict, None]:
    """
    Parse multiple files with streaming status updates.
    Files are processed in parallel with controlled concurrency.

    Args:
        files: List of (file_bytes, filename) tuples
        mode: The parsing mode (ignored for Docling)
        api_key: API key (ignored for Docling)

    Yields:
        Status updates and results as dicts
    """
    total = len(files)

    yield {"type": "start", "total": total, "mode": mode}

    if total == 0:
        yield {
            "type": "complete",
            "results": [],
            "success_count": 0,
            "error_count": 0
        }
        return

    # Queue for collecting progress events from parallel tasks
    event_queue: asyncio.Queue[dict] = asyncio.Queue()

    # Semaphore to limit concurrent parsing requests
    semaphore = asyncio.Semaphore(MAX_CONCURRENT_PARSE)

    # Results storage (index -> result)
    results: dict[int, dict] = {}

    async def parse_single_file(index: int, file_bytes: bytes, filename: str):
        """Parse a single file and put progress events into the queue."""
        async with semaphore:
            # Emit "parsing" status
            await event_queue.put({
                "type": "progress",
                "current": index + 1,
                "total": total,
                "filename": filename,
                "status": "parsing"
            })

            try:
                if is_simple_text(filename):
                    # Text file - read directly
                    await event_queue.put({
                        "type": "progress",
                        "current": index + 1,
                        "total": total,
                        "filename": filename,
                        "status": "reading_text"
                    })
                    try:
                        content = file_bytes.decode('utf-8')
                    except UnicodeDecodeError:
                        content = file_bytes.decode('latin-1')

                    results[index] = {
                        "filename": filename,
                        "content": content,
                        "parsed": False,
                        "error": None
                    }

                elif needs_parsing(filename):
                    # Complex file - use Docling
                    await event_queue.put({
                        "type": "progress",
                        "current": index + 1,
                        "total": total,
                        "filename": filename,
                        "status": "docling"
                    })

                    content = await parse_file(file_bytes, filename, mode, api_key=api_key)
                    results[index] = {
                        "filename": filename,
                        "content": content,
                        "parsed": True,
                        "error": None
                    }

                else:
                    # Unsupported file type
                    ext = Path(filename).suffix.lower()
                    results[index] = {
                        "filename": filename,
                        "content": None,
                        "parsed": False,
                        "error": f"Unsupported file type: {ext}"
                    }

                # Emit completion status
                await event_queue.put({
                    "type": "progress",
                    "current": index + 1,
                    "total": total,
                    "filename": filename,
                    "status": "complete"
                })

            except Exception as e:
                results[index] = {
                    "filename": filename,
                    "content": None,
                    "parsed": False,
                    "error": str(e)
                }
                await event_queue.put({
                    "type": "progress",
                    "current": index + 1,
                    "total": total,
                    "filename": filename,
                    "status": "error",
                    "error": str(e)
                })

    # Create tasks for all files
    tasks = [
        asyncio.create_task(parse_single_file(i, file_bytes, filename))
        for i, (file_bytes, filename) in enumerate(files)
    ]

    # Sentinel to signal all tasks are done
    async def wait_for_tasks():
        await asyncio.gather(*tasks)
        await event_queue.put({"type": "_done"})

    asyncio.create_task(wait_for_tasks())

    # Yield events as they arrive
    completed_count = 0
    while True:
        event = await event_queue.get()

        if event["type"] == "_done":
            break

        # Track completions for final summary
        if event["type"] == "progress" and event.get("status") in ("complete", "error"):
            completed_count += 1

        yield event

    # Build final results in original order
    ordered_results = [results[i] for i in range(total)]

    yield {
        "type": "complete",
        "results": ordered_results,
        "success_count": sum(1 for r in ordered_results if r["error"] is None),
        "error_count": sum(1 for r in ordered_results if r["error"] is not None)
    }


class ParsedFile:
    """Represents a parsed file with its content."""

    def __init__(
        self,
        filename: str,
        content: str,
        original_bytes: Optional[bytes] = None,
        was_parsed: bool = False
    ):
        self.filename = filename
        self.content = content
        self.original_bytes = original_bytes
        self.was_parsed = was_parsed

    def to_dict(self) -> dict:
        return {
            "filename": self.filename,
            "content": self.content,
            "was_parsed": self.was_parsed,
            # Don't include original_bytes in dict - it's for internal use
        }

    @classmethod
    def from_dict(cls, data: dict) -> "ParsedFile":
        return cls(
            filename=data["filename"],
            content=data["content"],
            was_parsed=data.get("was_parsed", False)
        )
